# -*- coding: utf-8 -*-
"""
Created on Tue Oct 15 18:15:36 2024

@author: Andres
"""

import numpy as np
import cv2
import matplotlib.pyplot as plt


font = cv2.FONT_HERSHEY_SIMPLEX 
  
org = (100, 100) 
fontScale = 1.5
color = (0, 0, 0)      
thickness = 8

from os import listdir
from skimage import segmentation
path = 'D:/TCGA MV Patches/MV_patches/'

files = listdir(path)

#%%
for file in files:
    seg  = cv2.imread('D:/TCGA MV Patches/MV_Mask/SG_'+file,cv2.IMREAD_GRAYSCALE)
    main = cv2.imread('D:/TCGA MV Patches/MV_patches/'+file,cv2.IMREAD_COLOR)
    

    result_image = segmentation.mark_boundaries(main, seg, mode='thick',color=(1,0,0),outline_color=(1,0,0))
    
    filename = file[:-4]
    
    zz = cv2.putText(result_image, filename, org, font,  
                       fontScale, color,thickness, cv2.LINE_AA) 
    
    # plt.figure()
    # plt.imshow(result_image)
    result_image = np.uint8(np.round(zz*255))
    cv2.imwrite('D:/X/'+'NSG_'+file,result_image) 

#%%

from PIL import Image

from pptx import Presentation
from pptx.util import Inches


# Open an image file
image = Image.open('D:/TCGA MV Patches/MV_patches/'+file)
# image.show()  # This will display the image

# Create a presentation object
# prs = Presentation()
#%%
# Add a slide with a title and content layout
# slide_layout = prs.slide_layouts  # Choosing a blank slide layout
presentation = Presentation()
slide_layout = presentation.slide_layouts[0]
# layout = presentation.slide_masters[0].slide_layouts[6]
slide = presentation.slides.add_slide(slide_layout)
# slide_id = presentation.slides.index(slide)
# slide = prs.slides.add_slide(presentation)

#Add the image to the slide
img_path = 'D:/X/NSG_TCGA-02-0001-01Z-00-DX1_000.png'
left = Inches(1)
top = Inches(0)
height = Inches(7.5)
slide.shapes.add_picture(img_path, left, top, height=height)

# Save the presentation
presentation.save('D:/Presentation1.pptx')

#%%

from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders
title.text = "Hello, World!"
subtitle.text = "This is a subtitle."

# Add a slide with a title and content layout
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders
title.text = "Slide Title"
content.text = "This is the content of the slide."

# Add a blank slide and an image
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
img_path = 'path_to_your_image.jpg'
left = Inches(1)
top = Inches(1)
height = Inches(5.5)
slide.shapes.add_picture(img_path, left, top, height=height)

# Save the presentation
prs.save('complete_example.pptx')

