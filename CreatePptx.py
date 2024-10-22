# -*- coding: utf-8 -*-
"""
Created on Tue Oct 15 18:51:48 2024

@author: Andres
"""

from PIL import Image

from pptx import Presentation
from pptx.util import Inches
from os import listdir


path = 'D:/TCGA MV Patches/MV_patches/'
files = listdir(path)

prs = Presentation()

for file in files[51:100]:
    image = Image.open('D:/TCGA MV Patches/MV_patches/'+file)
    print(file)
    
    
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    prs.slide_width = Inches(12)
    prs.slide_height = Inches(12)
    title = slide.shapes.title
    img_path = 'D:/TCGA MV Patches/MV_patches/'+file
    left = Inches(0)
    top = Inches(0)
    height = Inches(12)
    slide.shapes.add_picture(img_path, left, top, height=height)
    
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    img_path = 'D:/X/'+'NSG_'+file
    left = Inches(0)
    top = Inches(0)
    height = Inches(12)
    slide.shapes.add_picture(img_path, left, top, height=height)
    
    # subtitle = slide.placeholders
    # title.text = "Hello, World!"
    # subtitle.text = "This is a subtitle."
    
    # presentation = Presentation()
    # slide_layout = presentation.slide_masters[0].slide_layouts[6]
    # # layout = presentation.slide_masters[0].slide_layouts[6]
    # slide = presentation.slides.add_slide(slide_layout)
    # # slide_id = presentation.slides.index(slide)
    

    #Add the image to the slide
    # img_path = 'D:/TCGA MV Patches/MV_patches/'+file
    # left = Inches(1)
    # top = Inches(0)
    # height = Inches(7.5)
    # slide.shapes.add_picture(img_path, left, top, height=height)
    
prs.save('D:/MVP_51_to_100.pptx')
#%%
# Open an image file
# image = Image.open('D:/TCGA MV Patches/MV_patches/'+file)
# image.show()  # This will display the image

# Create a presentation object
# prs = Presentation()
#%%
# # Add a slide with a title and content layout
# # slide_layout = prs.slide_layouts  # Choosing a blank slide layout
# presentation = Presentation()
# slide_layout = presentation.slide_layouts[0]
# # layout = presentation.slide_masters[0].slide_layouts[6]
# slide = presentation.slides.add_slide(slide_layout)
# # slide_id = presentation.slides.index(slide)
# # slide = prs.slides.add_slide(presentation)

# #Add the image to the slide
# img_path = 'D:/X/NSG_TCGA-02-0001-01Z-00-DX1_000.png'
# left = Inches(1)
# top = Inches(0)
# height = Inches(7.5)
# slide.shapes.add_picture(img_path, left, top, height=height)

# # Save the presentation
# presentation.save('D:/Presentation1.pptx')

# #%%

# from pptx import Presentation
# from pptx.util import Inches

# # Create a presentation object
# prs = Presentation()

# # Add a title slide
# slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders
# title.text = "Hello, World!"
# subtitle.text = "This is a subtitle."

# # Add a slide with a title and content layout
# slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(slide_layout)
# title = slide.shapes.title
# content = slide.placeholders
# title.text = "Slide Title"
# content.text = "This is the content of the slide."

# # Add a blank slide and an image
# # slide_layout = prs.slide_layouts[0]
# # slide = prs.slides.add_slide(slide_layout)
# # img_path = 'path_to_your_image.jpg'
# # left = Inches(1)
# # top = Inches(1)
# # height = Inches(5.5)
# # slide.shapes.add_picture(img_path, left, top, height=height)

# # Save the presentation
# prs.save('D:/Presentation1.pptx')
